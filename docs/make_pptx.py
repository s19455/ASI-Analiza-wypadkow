"""Generuje prezentacje PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def add_slide(prs, layout_idx=1):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = RGBColor(*color)
    return tf


def add_para(tf, text, size=18, bold=False, color=None, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.space_before = space_before
    if color:
        p.font.color.rgb = RGBColor(*color)
    return p


def emu(inches):
    return Inches(inches)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height
    MARGIN = Inches(0.8)
    WHITE = (255, 255, 255)
    DARK = (44, 62, 80)
    BLUE = (52, 152, 219)
    RED = (231, 76, 60)
    GREEN = (39, 174, 96)
    GRAY = (127, 140, 141)

    # ===== SLAJD 1: Tytul =====
    s = add_slide(prs, 6)  # blank
    set_bg(s, 44, 62, 80)
    add_textbox(s, MARGIN, Inches(1.8), W - 2*MARGIN, Inches(1.5),
        "Analiza wypadków drogowych", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, MARGIN, Inches(3.3), W - 2*MARGIN, Inches(0.8),
        "Predykcja stopnia obrażeń", size=28, color=(170, 170, 170), align=PP_ALIGN.CENTER)
    add_textbox(s, MARGIN, Inches(4.5), W - 2*MARGIN, Inches(0.6),
        "Artur Cichocki  •  Bartosz Pikutin  •  Wiktor Golba", size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, MARGIN, Inches(5.3), W - 2*MARGIN, Inches(0.5),
        "Projekt zaliczeniowy ASI — PJATK", size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(s, MARGIN, Inches(5.8), W - 2*MARGIN, Inches(0.5),
        "Dane: Montgomery County, Maryland (2015–2024)", size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # ===== SLAJD 2: Problem =====
    s = add_slide(prs, 6)
    set_bg(s, 255, 255, 255)
    add_textbox(s, MARGIN, Inches(0.4), Inches(5), Inches(0.7),
        "Problem", size=36, bold=True, color=DARK)

    tf = add_textbox(s, MARGIN, Inches(1.3), Inches(6), Inches(4),
        "Cel: na podstawie warunków wypadku przewidzieć stopień obrażeń.", size=20, color=DARK)
    add_para(tf, "", size=10)
    add_para(tf, "Dlaczego to ważne:", size=20, bold=True, color=DARK)
    add_para(tf, "•  Planowanie bezpieczeństwa ruchu", size=18, color=DARK)
    add_para(tf, "•  Lepsze dysponowanie służb ratunkowych", size=18, color=DARK)
    add_para(tf, "•  Zrozumienie czynników ryzyka", size=18, color=DARK)

    tf2 = add_textbox(s, Inches(7.5), Inches(1.3), Inches(5), Inches(4),
        "3 klasy:", size=20, bold=True, color=DARK)
    add_para(tf2, "", size=10)
    add_para(tf2, "NO_INJURY — brak obrażeń (82%)", size=20, bold=True, color=GREEN)
    add_para(tf2, "MINOR — drobne obrażenia (17%)", size=20, bold=True, color=(230, 160, 0))
    add_para(tf2, "SERIOUS — poważne / zgon (1%)", size=20, bold=True, color=RED)
    add_para(tf2, "", size=10)
    add_para(tf2, "Problem: silna nierównowaga klas", size=16, color=GRAY)

    # ===== SLAJD 3: Dane =====
    s = add_slide(prs, 6)
    set_bg(s, 255, 255, 255)
    add_textbox(s, MARGIN, Inches(0.4), Inches(5), Inches(0.7),
        "Zbiór danych", size=36, bold=True, color=DARK)

    for i, (val, label) in enumerate([("172 105", "rekordów"), ("43", "kolumny"), ("10 lat", "danych")]):
        x = Inches(1.5 + i * 3.5)
        add_textbox(s, x, Inches(1.5), Inches(3), Inches(0.8),
            val, size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(2.3), Inches(3), Inches(0.5),
            label, size=18, color=GRAY, align=PP_ALIGN.CENTER)

    tf = add_textbox(s, MARGIN, Inches(3.3), Inches(6), Inches(3.5),
        "Ważniejsze cechy:", size=20, bold=True, color=DARK)
    add_para(tf, "•  Pogoda, oświetlenie, stan nawierzchni", size=17, color=DARK)
    add_para(tf, "•  Typ kolizji, ruch pojazdu", size=17, color=DARK)
    add_para(tf, "•  Limit prędkości, typ pojazdu", size=17, color=DARK)
    add_para(tf, "•  Data, godzina, lokalizacja GPS", size=17, color=DARK)

    tf2 = add_textbox(s, Inches(7.5), Inches(3.3), Inches(5), Inches(3.5),
        "Przetwarzanie:", size=20, bold=True, color=DARK)
    add_para(tf2, "•  Usunięcie kolumn z >90% braków", size=17, color=DARK)
    add_para(tf2, "•  Imputacja medianą / modą", size=17, color=DARK)
    add_para(tf2, "•  Nowe cechy: is_night, vehicle_age...", size=17, color=DARK)
    add_para(tf2, "•  Label encoding (wysoka kardynalność)", size=17, color=DARK)

    # ===== SLAJD 4: EDA =====
    s = add_slide(prs, 6)
    set_bg(s, 255, 255, 255)
    add_textbox(s, MARGIN, Inches(0.4), Inches(8), Inches(0.7),
        "EDA — co zauważyliśmy", size=36, bold=True, color=DARK)

    tf = add_textbox(s, MARGIN, Inches(1.3), Inches(6), Inches(5),
        "Główne obserwacje:", size=20, bold=True, color=DARK)
    add_para(tf, "•  Wypadki SERIOUS częstsze w nocy i przy złej pogodzie", size=18, color=DARK)
    add_para(tf, "•  Wyższy limit prędkości → więcej poważnych obrażeń", size=18, color=DARK)
    add_para(tf, "•  Kolizje czołowe — najwyższy % SERIOUS", size=18, color=DARK)
    add_para(tf, "•  Szczyty wypadków: 7-9 i 16-18", size=18, color=DARK)
    add_para(tf, "•  Piątki — najwięcej wypadków", size=18, color=DARK)

    tf2 = add_textbox(s, Inches(7.5), Inches(1.3), Inches(5), Inches(5),
        "Top cechy (wg Random Forest):", size=20, bold=True, color=DARK)
    add_para(tf2, "1.  Vehicle Damage Extent", size=18, color=DARK)
    add_para(tf2, "2.  Collision Type", size=18, color=DARK)
    add_para(tf2, "3.  Speed Limit", size=18, color=DARK)
    add_para(tf2, "4.  Vehicle Movement", size=18, color=DARK)
    add_para(tf2, "5.  Crash hour", size=18, color=DARK)
    add_para(tf2, "", size=10)
    add_para(tf2, "Szczegóły: notebooks/01_baseline_model.ipynb", size=14, color=GRAY)

    # ===== SLAJD 5: Pipeline =====
    s = add_slide(prs, 6)
    set_bg(s, 255, 255, 255)
    add_textbox(s, MARGIN, Inches(0.4), Inches(8), Inches(0.7),
        "Pipeline ML (Kedro)", size=36, bold=True, color=DARK)

    pipeline_text = (
        "crash_data.csv\n"
        "      ↓\n"
        "[data_preparation] — 5 węzłów\n"
        "  usunięcie kolumn → imputacja → inżynieria cech\n"
        "  → mapowanie target → enkodowanie\n"
        "      ↓\n"
        "[data_modeling] — 3 węzły\n"
        "  podział train/test → trening → ewaluacja\n"
        "      ↓\n"
        "[tuning]\n"
        "  porównanie 4 modeli + Optuna (50 prób)"
    )
    add_textbox(s, Inches(1.5), Inches(1.5), Inches(10), Inches(4.5),
        pipeline_text, size=20, color=DARK)

    add_textbox(s, MARGIN, Inches(6.2), Inches(10), Inches(0.5),
        "Uruchomienie:  kedro run  |  kedro run --pipeline=tuning  |  kedro viz", size=16, color=GRAY)

    # ===== SLAJD 6: Wyniki =====
    s = add_slide(prs, 6)
    set_bg(s, 255, 255, 255)
    add_textbox(s, MARGIN, Inches(0.4), Inches(8), Inches(0.7),
        "Wyniki modeli", size=36, bold=True, color=DARK)

    # tabela
    rows, cols = 5, 4
    tbl = s.shapes.add_table(rows, cols, Inches(1), Inches(1.4), Inches(11), Inches(2.8)).table

    headers = ["Model", "Dokładność", "F1 ważone", "F1 makro"]
    data = [
        ["Optuna LightGBM", "0.78", "0.78", "0.46"],
        ["XGBoost", "0.83", "0.77", "0.40"],
        ["Gradient Boosting", "0.82", "0.77", "0.40"],
        ["Random Forest (baseline)", "0.82", "0.75", "0.33"],
    ]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*BLUE)

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(15)
                if i == 0:
                    p.font.bold = True

    # podswietlenie najlepszego
    for j in range(cols):
        tbl.cell(1, j).fill.solid()
        tbl.cell(1, j).fill.fore_color.rgb = RGBColor(212, 239, 223)

    tf = add_textbox(s, MARGIN, Inches(4.5), Inches(11), Inches(1.5),
        "F1 makro jest tu ważniejsze niż accuracy — mierzymy jak model radzi sobie z rzadkimi klasami.", size=18, color=DARK)
    add_para(tf, "Poprawa F1 makro: 0.33 → 0.46 (+39% vs baseline)", size=18, bold=True, color=DARK)
    add_para(tf, "Optuna: 50 prób optymalizacji bayesowskiej", size=14, color=GRAY)

    # ===== SLAJD 7: Wdrozenie =====
    s = add_slide(prs, 6)
    set_bg(s, 255, 255, 255)
    add_textbox(s, MARGIN, Inches(0.4), Inches(8), Inches(0.7),
        "Wdrożenie", size=36, bold=True, color=DARK)

    tf = add_textbox(s, MARGIN, Inches(1.3), Inches(5.5), Inches(5),
        "API (FastAPI):", size=22, bold=True, color=DARK)
    add_para(tf, "•  POST /predict — predykcja", size=17, color=DARK)
    add_para(tf, "•  GET /health — status", size=17, color=DARK)
    add_para(tf, "•  Swagger UI pod /docs", size=17, color=DARK)
    add_para(tf, "", size=10)
    add_para(tf, "Docker:", size=22, bold=True, color=DARK)
    add_para(tf, "•  Dockerfile + docker-compose", size=17, color=DARK)
    add_para(tf, "•  docker-compose up --build", size=17, color=DARK)

    tf2 = add_textbox(s, Inches(7), Inches(1.3), Inches(5.5), Inches(5),
        "CI/CD (GitHub Actions):", size=22, bold=True, color=DARK)
    add_para(tf2, "•  CI — ruff (lint) + pytest", size=17, color=DARK)
    add_para(tf2, "•  CD — build Docker → push do GHCR", size=17, color=DARK)
    add_para(tf2, "•  CT — automatyczne retrenowanie", size=17, color=DARK)
    add_para(tf2, "", size=10)
    add_para(tf2, "Monitoring:", size=22, bold=True, color=DARK)
    add_para(tf2, "•  Logowanie każdej predykcji", size=17, color=DARK)
    add_para(tf2, "•  Evidently — wykrywanie driftu", size=17, color=DARK)

    # ===== SLAJD 8: Technologie =====
    s = add_slide(prs, 6)
    set_bg(s, 255, 255, 255)
    add_textbox(s, MARGIN, Inches(0.4), Inches(8), Inches(0.7),
        "Użyte technologie", size=36, bold=True, color=DARK)

    tf = add_textbox(s, Inches(1.5), Inches(1.5), Inches(5), Inches(5),
        "", size=18, color=DARK)
    for item in ["Kedro — pipeline ML", "scikit-learn — modele bazowe", "XGBoost, LightGBM — gradient boosting",
                 "Optuna — strojenie hiperparametrów", "MLflow — śledzenie eksperymentów"]:
        add_para(tf, "•  " + item, size=18, color=DARK)

    tf2 = add_textbox(s, Inches(7), Inches(1.5), Inches(5), Inches(5),
        "", size=18, color=DARK)
    for item in ["FastAPI — API do predykcji", "Docker — konteneryzacja", "GitHub Actions — CI/CD",
                 "Evidently — monitoring driftu", "Jupyter — EDA i eksperymenty"]:
        add_para(tf2, "•  " + item, size=18, color=DARK)

    add_textbox(s, MARGIN, Inches(5.5), Inches(11), Inches(0.5),
        "Repo: github.com/s19455/ASI-Analiza-wypadkow", size=16, color=GRAY)

    # ===== SLAJD 9: Podsumowanie =====
    s = add_slide(prs, 6)
    set_bg(s, 44, 62, 80)

    add_textbox(s, MARGIN, Inches(0.8), W - 2*MARGIN, Inches(0.8),
        "Podsumowanie", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (val, label) in enumerate([("172K", "wypadków"), ("4", "modele"), ("0.78", "F1 ważone"), ("+39%", "poprawa F1 makro")]):
        x = Inches(1.2 + i * 3)
        add_textbox(s, x, Inches(2.3), Inches(2.8), Inches(0.8),
            val, size=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_textbox(s, x, Inches(3.1), Inches(2.8), Inches(0.5),
            label, size=18, color=GRAY, align=PP_ALIGN.CENTER)

    add_textbox(s, MARGIN, Inches(4.5), W - 2*MARGIN, Inches(0.6),
        "Od surowych danych, przez pipeline Kedro, po API w Dockerze z CI/CD",
        size=20, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, MARGIN, Inches(5.8), W - 2*MARGIN, Inches(0.8),
        "Dziękujemy — pytania?", size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Zapis
    prs.save("docs/prezentacja.pptx")
    print("Zapisano docs/prezentacja.pptx")


if __name__ == "__main__":
    main()
